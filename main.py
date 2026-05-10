from flask import Flask, request, jsonify
import requests
import re
from datetime import datetime, timezone, timedelta
import json
import os
import threading
import time
import base64
import io

# 文件解析库
import PyPDF2
import docx
import openpyxl
import pptx

app = Flask(__name__)

# ========== 配置 ==========
SLACK_BOT_TOKEN = os.environ.get("SLACK_BOT_TOKEN")
JSONBIN_API_KEY = os.environ.get("JSONBIN_API_KEY")
JSONBIN_USER_DATA = os.environ.get("JSONBIN_USER_DATA")
JSONBIN_SCHEDULES = os.environ.get("JSONBIN_SCHEDULES")
JSONBIN_MEMORIES = os.environ.get("JSONBIN_MEMORIES")
JSONBIN_CHAT_LOGS = os.environ.get("JSONBIN_CHAT_LOGS")
JSONBIN_CHANNEL_MESSAGES = os.environ.get("JSONBIN_CHANNEL_MESSAGES")

API_TOKEN_LIMITS = {
    "第三方sonnet": 110000,
    "sonnet": 190000,
    "gptsapi-sonnet": 190000,
    "gptsapi-thinking": 190000,
    "gptsapi-opus": 190000,
    "gptsapi-opus4.7": 190000,
    "gptsapi-haiku": 190000,
    "code haiku": 190000,
    "code sonnet": 190000,
    "code opus": 190000,
    "啾啾sonnet": 190000,
    "啾啾opus": 190000,
    "R Opus": 190000,
    "小E Opus thinking": 190000,
    "小E Opus": 190000,
}

APIS = {
    # === 默认API (小鸡农场) ===
    "sonnet": {
        "url": os.environ.get("API_URL_1"),
        "key": os.environ.get("API_KEY_1"),
        "model": "claude-sonnet-4-5 [官转1]",
        "vision": True,
        "cost": 4
    },
    # === 新API (gptsapi) ===
    "gptsapi-sonnet": {
        "url": "https://api.gptsapi.net/v1/chat/completions",
        "key": "sk-H4va4c84c4791ef6216efba7cd161d478916efb2c45ICCXz",
        "model": "claude-sonnet-4-6",
        "vision": True,
        "cost": 4
    },
    "gptsapi-thinking": {
        "url": "https://api.gptsapi.net/v1/chat/completions",
        "key": "sk-H4va4c84c4791ef6216efba7cd161d478916efb2c45ICCXz",
        "model": "claude-sonnet-4-6-thinking",
        "vision": True,
        "cost": 5
    },
    "gptsapi-opus": {
        "url": "https://api.gptsapi.net/v1/chat/completions",
        "key": "sk-H4va4c84c4791ef6216efba7cd161d478916efb2c45ICCXz",
        "model": "claude-opus-4-6",
        "vision": True,
        "cost": 10
    },
    "gptsapi-opus4.7": {
        "url": "https://api.gptsapi.net/v1/chat/completions",
        "key": "sk-H4va4c84c4791ef6216efba7cd161d478916efb2c45ICCXz",
        "model": "claude-opus-4-7",
        "vision": True,
        "cost": 10
    },
    "gptsapi-haiku": {
        "url": "https://api.gptsapi.net/v1/chat/completions",
        "key": "sk-H4va4c84c4791ef6216efba7cd161d478916efb2c45ICCXz",
        "model": "claude-haiku-4-5-20251001",
        "vision": True,
        "cost": 2
    },
    # === 旧API (通过环境变量配置) ===
    "第三方sonnet": {
        "url": os.environ.get("API_URL_1"),
        "key": os.environ.get("API_KEY_1"),
        "model": "[第三方逆1] claude-sonnet-4.5 [输出只有3~4k]",
        "vision": False,
        "cost": 1
    },
    "code haiku": {
        "url": os.environ.get("API_URL_3"),
        "key": os.environ.get("API_KEY_3"),
        "model": "[code]claude-haiku-4-5-20251001",
        "vision": True,
        "cost": 2
    },
    "code sonnet": {
        "url": os.environ.get("API_URL_3"),
        "key": os.environ.get("API_KEY_3"),
        "model": "[code]claude-sonnet-4-5-20250929",
        "vision": True,
        "cost": 5
    },
    "code opus": {
        "url": os.environ.get("API_URL_3"),
        "key": os.environ.get("API_KEY_3"),
        "model": "[code]claude-opus-4-5-20251101",
        "vision": True,
        "cost": 10
    },
    "啾啾sonnet": {
        "url": os.environ.get("API_URL_3"),
        "key": os.environ.get("API_KEY_3"),
        "model": "[啾啾]claude-sonnet-4-5-20250929",
        "vision": True,
        "cost": 5
    },
    "啾啾opus": {
        "url": os.environ.get("API_URL_3"),
        "key": os.environ.get("API_KEY_3"),
        "model": "[啾啾]claude-opus-4-5-20251101",
        "vision": True,
        "cost": 10
    },
    "R Opus": {
        "url": os.environ.get("API_URL_2"),
        "key": os.environ.get("API_KEY_2"),
        "model": "R-claude-opus-4-6",
        "vision": True,
        "cost": 10
    },
    "小E Opus thinking": {
        "url": os.environ.get("API_URL_2"),
        "key": os.environ.get("API_KEY_2"),
        "model": "小E-claude-opus-4-6-thinking",
        "vision": True,
        "cost": 10
    },
    "小E Opus": {
        "url": os.environ.get("API_URL_2"),
        "key": os.environ.get("API_KEY_2"),
        "model": "小E-claude-opus-4-6",
        "vision": True,
        "cost": 10
    },
}

DEFAULT_API = "小E Opus"  # gptsapi余额耗尽，切到小E
UNLIMITED_USERS = ["sakuragochyan", "ms-sydney"]
POINTS_LIMIT = 20
MEMORY_LIMIT = 2000
CONVERSATION_TIMEOUT = 300
MAX_FILE_SIZE = 10 * 1024 * 1024

# AI 积分系统配置
AI_POINTS_MAX = 10
AI_POINTS_MIN = -10
AI_POINTS_DEFAULT = 10
AI_MSG_LENGTH_LIMIT = 50  # 超过就审查
AI_MSG_LENGTH_IDEAL = 20  # 理想长度
REVIEW_TOKEN_LIMIT = 110000  # 审查时的 token 限制
MAX_REWORK_ATTEMPTS = 3  # 最大返工次数

CN_TIMEZONE = timezone(timedelta(hours=8))

processed_events = set()
processed_file_events = set()
pending_messages = {}
pending_timers = {}
pending_clear_logs = {}
channel_message_counts = {}

NO_DM_HISTORY_CHANNELS = ["learn"]

EMOJI_ALIASES = {
    "thumbs_up": "thumbsup", "thumb_up": "thumbsup", "+1": "thumbsup", "like": "thumbsup",
    "thumbs_down": "thumbsdown", "thumb_down": "thumbsdown", "-1": "thumbsdown", "dislike": "thumbsdown",
    "joy": "laughing", "sob": "cry", "crying": "cry", "sad": "cry",
    "love": "heart", "red_heart": "heart",
    "think": "thinking_face", "thinking": "thinking_face", "hmm": "thinking_face",
    "clapping": "clap", "applause": "clap",
    "party": "tada", "celebrate": "tada", "celebration": "tada",
    "stars": "sparkles", "glitter": "sparkles", "shine": "sparkles",
    "hi": "wave", "hello": "wave", "bye": "wave",
    "thanks": "pray", "thank_you": "pray", "please": "pray", "gratitude": "pray",
    "hundred": "100", "perfect": "100",
    "flame": "fire", "hot": "fire", "lit": "fire",
    "look": "eyes", "see": "eyes", "watching": "eyes",
    "ok": "ok_hand", "okay": "ok_hand",
    "strong": "muscle", "strength": "muscle", "flex": "muscle",
    "cool": "sunglasses",
    "check": "white_check_mark", "yes": "white_check_mark",
    "no": "x", "wrong": "x",
    "sleep": "zzz", "sleepy": "zzz", "tired": "zzz",
    "sweat": "sweat_smile", "nervous": "sweat_smile",
}

VALID_EMOJIS = [
    "heart", "thumbsup", "thumbsdown", "laughing", "cry", "fire", 
    "eyes", "thinking_face", "clap", "tada", "star", "wave", 
    "pray", "sparkles", "100", "rocket", "muscle", "ok_hand", 
    "raised_hands", "sunglasses", "white_check_mark", "x", "zzz", 
    "sweat_smile", "blush", "wink", "grin", "smile"
]

TEXT_EXTENSIONS = ['.txt', '.md', '.py', '.js', '.ts', '.html', '.css', '.json', '.xml', '.yaml', '.yml', '.csv', '.log', '.sh', '.bash', '.c', '.cpp', '.h', '.java', '.rb', '.php', '.go', '.rs', '.swift', '.kt', '.r', '.sql']

# ========== 文件解析函数 ==========

def download_file(url):
    try:
        resp = requests.get(url, headers={"Authorization": f"Bearer {SLACK_BOT_TOKEN}"}, timeout=30)
        if resp.status_code == 200:
            return resp.content
    except Exception as e:
        print(f"[File] 下载失败: {e}")
    return None

def extract_pdf_text(content):
    try:
        pdf = PyPDF2.PdfReader(io.BytesIO(content))
        text = ""
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
        return text.strip() if text.strip() else None
    except Exception as e:
        print(f"[File] PDF 解析失败: {e}")
        return None

def extract_docx_text(content):
    try:
        doc = docx.Document(io.BytesIO(content))
        return "\n".join([para.text for para in doc.paragraphs]).strip() or None
    except Exception as e:
        print(f"[File] Word 解析失败: {e}")
        return None

def extract_xlsx_text(content):
    try:
        wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True)
        text = ""
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            text += f"[工作表: {sheet}]\n"
            for row in ws.iter_rows(values_only=True):
                row_text = "\t".join([str(cell) if cell else "" for cell in row])
                if row_text.strip():
                    text += row_text + "\n"
        return text.strip() or None
    except Exception as e:
        print(f"[File] Excel 解析失败: {e}")
        return None

def extract_pptx_text(content):
    try:
        prs = pptx.Presentation(io.BytesIO(content))
        text = ""
        for i, slide in enumerate(prs.slides, 1):
            text += f"[幻灯片 {i}]\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
        return text.strip() or None
    except Exception as e:
        print(f"[File] PPT 解析失败: {e}")
        return None

def extract_text_file(content):
    for encoding in ['utf-8', 'gbk', 'gb2312', 'latin-1']:
        try:
            return content.decode(encoding)
        except:
            continue
    return None

def process_file(file_info):
    filename = file_info.get("name", "未知文件")
    mimetype = file_info.get("mimetype", "")
    file_size = file_info.get("size", 0)
    url = file_info.get("url_private")
    
    if not url:
        return None, None
    if file_size > MAX_FILE_SIZE:
        return "too_large", f"[文件: {filename}]（超过 10MB 限制）"
    
    ext = os.path.splitext(filename)[1].lower()
    
    if mimetype.startswith("image/"):
        return "image", url
    
    content = download_file(url)
    if not content:
        return "error", f"[文件: {filename}]（下载失败）"
    
    parsers = {
        ".pdf": extract_pdf_text,
        ".docx": extract_docx_text,
        ".xlsx": extract_xlsx_text,
        ".pptx": extract_pptx_text,
    }
    
    if ext in parsers:
        text = parsers[ext](content)
        return ("text", f"[文件: {filename}]\n{text}") if text else ("error", f"[文件: {filename}]（解析失败）")
    
    if ext in [".doc", ".xls", ".ppt"]:
        return "unsupported", f"[文件: {filename}]（不支持旧版格式）"
    
    if ext in TEXT_EXTENSIONS or mimetype.startswith("text/"):
        text = extract_text_file(content)
        if text:
            if len(text) > 50000:
                text = text[:50000] + "\n...(已截断)"
            return "text", f"[文件: {filename}]\n{text}"
    
    return "unsupported", f"[文件: {filename}]（不支持此格式）"

# ========== JSONBin 工具 ==========

def jsonbin_save(bin_id, data):
    if not bin_id:
        return
    try:
        requests.put(
            f"https://api.jsonbin.io/v3/b/{bin_id}",
            headers={"X-Master-Key": JSONBIN_API_KEY, "Content-Type": "application/json"},
            json=data, timeout=30
        )
    except Exception as e:
        print(f"JSONBin 保存失败: {e}")

def jsonbin_load(bin_id, default=None):
    if not bin_id:
        return default or {}
    try:
        resp = requests.get(
            f"https://api.jsonbin.io/v3/b/{bin_id}/latest",
            headers={"X-Master-Key": JSONBIN_API_KEY}, timeout=30
        )
        print(f"[JSONBin] load {bin_id}: status={resp.status_code}")
        if resp.status_code == 200:
            record = resp.json().get("record", default or {})
            record.pop("init", None)
            return record
        else:
            print(f"[JSONBin] load 失败: {resp.text[:200]}")
    except Exception as e:
        print(f"[JSONBin] load 出错: {e}")
    return default or {}

# ========== 时间工具 ==========

def get_cn_time():
    return datetime.now(CN_TIMEZONE)

def get_time_str():
    weekdays = ["一", "二", "三", "四", "五", "六", "日"]
    now = get_cn_time()
    return now.strftime("%Y年%m月%d日 %H:%M:%S 星期") + weekdays[now.weekday()]

def get_timestamp():
    return get_cn_time().strftime("%Y-%m-%d %H:%M:%S")

def get_time_period():
    hour = get_cn_time().hour
    if 5 <= hour < 9:
        return "早上", "早上好"
    elif 9 <= hour < 12:
        return "上午", "上午好"
    elif 12 <= hour < 14:
        return "中午", "吃饭了吗"
    elif 14 <= hour < 18:
        return "下午", "下午好"
    elif 18 <= hour < 22:
        return "晚上", "晚上好"
    else:
        return "深夜", "注意休息"

# ========== 数据持久化 ==========

def load_user_data():
    return jsonbin_load(JSONBIN_USER_DATA, {})

def save_user_data(data):
    jsonbin_save(JSONBIN_USER_DATA, data)

def load_schedules():
    return jsonbin_load(JSONBIN_SCHEDULES, {})

def save_schedules(data):
    jsonbin_save(JSONBIN_SCHEDULES, data)

def load_channel_messages():
    return jsonbin_load(JSONBIN_CHANNEL_MESSAGES, {})

def save_channel_messages(data):
    jsonbin_save(JSONBIN_CHANNEL_MESSAGES, data)

def add_channel_message(channel_id, user_id, username, content, is_bot=False):
    try:
        messages = load_channel_messages()
        if channel_id not in messages:
            messages[channel_id] = []
        
        messages[channel_id].append({
            "user_id": user_id,
            "username": username,
            "content": content,
            "timestamp": get_cn_time().timestamp(),
            "time_str": get_timestamp(),
            "is_bot": is_bot
        })
        
        if len(messages[channel_id]) > 200:
            messages[channel_id] = messages[channel_id][-200:]
        
        save_channel_messages(messages)
        return len([m for m in messages[channel_id] if not m.get("is_bot")])
    except Exception as e:
        print(f"add_channel_message 出错: {e}")
        return 0

def get_channel_messages_since_reset(channel_id, reset_time=None):
    try:
        messages = load_channel_messages()
        msgs = messages.get(channel_id, [])
        if reset_time:
            msgs = [m for m in msgs if m.get("timestamp", 0) > reset_time]
        return msgs
    except:
        return []

def get_recent_channel_messages(channel_id, count=10):
    try:
        messages = load_channel_messages()
        return messages.get(channel_id, [])[-count:]
    except:
        return []

# ========== 聊天记录 ==========

def load_chat_logs():
    return jsonbin_load(JSONBIN_CHAT_LOGS, {})

def save_chat_logs(data):
    jsonbin_save(JSONBIN_CHAT_LOGS, data)

def log_message(user_id, channel, role, content, username=None, model=None, is_reset=False, hidden=False):
    try:
        logs = load_chat_logs()
        if user_id not in logs:
            logs[user_id] = []
        
        timestamp = get_timestamp()
        scene = "私聊" if is_dm_channel(channel) else get_channel_name(channel)
        
        if is_reset:
            logs[user_id].append({"type": "reset", "time": timestamp, "scene": scene})
        else:
            entry = {"time": timestamp, "scene": scene, "role": role, "content": content, "hidden": hidden}
            if role == "user":
                entry["username"] = username or "未知"
            else:
                entry["model"] = model or "未知"
            logs[user_id].append(entry)
        
        logs[user_id] = sorted(logs[user_id], key=lambda x: x.get("time", ""))
        save_chat_logs(logs)
    except Exception as e:
        print(f"log_message 出错: {e}")

def clear_user_chat_logs(user_id, channel_only=None):
    try:
        logs = load_chat_logs()
        if user_id not in logs:
            return
        if channel_only:
            channel_name = get_channel_name(channel_only)
            logs[user_id] = [e for e in logs[user_id] if e.get("scene") != channel_name]
        else:
            logs[user_id] = []
        save_chat_logs(logs)
    except Exception as e:
        print(f"clear_user_chat_logs 出错: {e}")

# ========== AI 积分系统 ==========

def get_ai_points(user_id):
    all_data = load_user_data()
    return all_data.get(user_id, {}).get("ai_points", AI_POINTS_DEFAULT)

def set_ai_points(user_id, points):
    points = max(AI_POINTS_MIN, min(AI_POINTS_MAX, points))
    all_data = load_user_data()
    if user_id not in all_data:
        all_data[user_id] = {}
    old_points = all_data[user_id].get("ai_points", AI_POINTS_DEFAULT)
    all_data[user_id]["ai_points"] = points
    save_user_data(all_data)
    return old_points, points

def deduct_ai_points(user_id, reason=""):
    current = get_ai_points(user_id)
    if current <= AI_POINTS_MIN:
        print(f"[AI积分] 用户 {user_id} 已经是最低分 {AI_POINTS_MIN}，无法再扣")
        return current, current, True  # 返回 True 表示需要返工
    
    deduct = 2 if current > 0 else 5
    old, new = set_ai_points(user_id, current - deduct)
    print(f"[AI积分] 用户 {user_id} 扣分: {old} -> {new} (扣{deduct}), 原因: {reason}")
    return old, new, False

def reward_ai_points(user_id):
    current = get_ai_points(user_id)
    if current < AI_POINTS_MAX:
        old, new = set_ai_points(user_id, current + 1)
        print(f"[AI积分] 用户 {user_id} 加分: {old} -> {new}")
        return old, new
    return current, current

def get_ai_points_status(user_id):
    """获取积分状态和对应的提示信息"""
    points = get_ai_points(user_id)
    
    if points <= AI_POINTS_MIN:
        return points, "min", f"""
🔥🔥🔥 *你的积分已经是最低分 {AI_POINTS_MIN} 了！！！* 🔥🔥🔥
如果你再犯错（分点列举、消息过长、回复过多），你的回复会被*强制返工*！
你必须重新生成回复，直到符合要求为止！

*你现在必须*：
- 每条消息控制在 20 字以内（最多不超过 50 字）
- 不要分点列举！
- 回复条数不要超过用户消息数的 3 倍！"""
    
    elif points < 0:
        return points, "negative", f"""
💀 *你的积分是负数了！*（当前: {points}/{AI_POINTS_MAX}）
现在每次犯错扣 5 分！再扣到 {AI_POINTS_MIN} 就要强制返工了！
控制回复长度和数量！不要分点列举！"""
    
    elif points == 0:
        return points, "zero", f"""
💀 *你的积分是 0！*（当前: 0/{AI_POINTS_MAX}）
再犯错就是负分了，负分后每次扣 5 分！"""
    
    elif points <= 2:
        return points, "danger", f"""
🚨 *严重警告！积分只剩 {points} 了！*（{points}/{AI_POINTS_MAX}）
控制回复！不要分点列举！"""
    
    elif points <= 6:
        return points, "warning", f"""
⚠️ *警告：积分 {points}*（{points}/{AI_POINTS_MAX}）
注意控制回复长度和数量。"""
    
    else:
        return points, "ok", f"当前积分: {points}/{AI_POINTS_MAX}"

def estimate_tokens(text):
    if not text:
        return 0
    chinese = len(re.findall(r'[\u4e00-\u9fff]', str(text)))
    other = len(str(text)) - chinese
    return int(chinese / 1.5 + other / 4)

def build_review_context(user, current_channel, user_message, ai_reply, msg_count):
    """构建用于审查的上下文，限制在 REVIEW_TOKEN_LIMIT 内"""
    context_parts = []
    
    # 收集所有历史消息
    all_messages = []
    
    # 私聊历史
    if "dm_history" in user: # Ensure dm_history exists
        for m in user.get("dm_history", []):
            if m.get("content"):
                all_messages.append({
                    "content": f"[私聊][{'用户' if m['role']=='user' else 'AI'}] {m['content']}",
                    "timestamp": m.get("timestamp", 0)
                })
    
    # 频道历史
    if not is_dm_channel(current_channel):
        reset_time = user.get("channel_reset_times", {}).get(current_channel, 0)
        channel_msgs = get_channel_messages_since_reset(current_channel, reset_time)
        for m in channel_msgs:
            sender = "AI" if m.get("is_bot") else m.get("username", "某人")
            all_messages.append({
                "content": f"[频道][{sender}] {m.get('content', '')}",
                "timestamp": m.get("timestamp", 0)
            })
    
    # 按时间排序
    all_messages.sort(key=lambda x: x["timestamp"])
    
    # 计算 token 并删除旧消息
    total_tokens = 0
    for m in all_messages:
        total_tokens += estimate_tokens(m["content"])
    
    while total_tokens > REVIEW_TOKEN_LIMIT and all_messages:
        removed = all_messages.pop(0)
        total_tokens -= estimate_tokens(removed["content"])
    
    context = "\n".join([m["content"] for m in all_messages])
    return context

def check_reply_format_violation(reply):
    """检查分点列举（直接扣分不审查）"""
    if re.search(r'^\s*\d+\.\s', reply, re.MULTILINE):
        return True, "数字列表 (1. 2. 3.)"
    if re.search(r'^\s*[•·]\s', reply, re.MULTILINE):
        return True, "圆点列表 (•)"
    if re.search(r'^\s*-\s+\S', reply, re.MULTILINE):
        return True, "横线列表 (-)"
    return False, None

def check_messages_too_long(messages):
    """检查是否有超过 50 字的消息"""
    for msg in messages:
        msg = msg.strip()
        if len(msg) > AI_MSG_LENGTH_LIMIT:
            return True, msg, len(msg)
    return False, None, 0

def review_with_ai(user, current_channel, user_message, ai_reply, msg_count, issue_type, details=""):
    """用第三方 sonnet 审查，带完整上下文"""
    try:
        api = APIS["第三方sonnet"]
        
        # 构建上下文
        context = build_review_context(user, current_channel, user_message, ai_reply, msg_count)
        
        if issue_type == "count":
            reply_count = len([m.strip() for m in ai_reply.split("|||") if m.strip()]) if "|||" in ai_reply else 1
            prompt = f"""你是一个审查员，需要判断 AI 的回复是否合理。

=== 聊天记录 ===
{context}

=== 当前情况 ===
用户发了 {msg_count} 条消息：
{user_message}

AI 回复了 {reply_count} 条消息。

=== 判断标准 ===
"不合理"的定义：AI 回复条数超过用户消息数的 3 倍（即超过 {msg_count * 3} 条），且用户的问题并不复杂，也没有要求 AI 回复多条。

请根据聊天记录和当前情况，判断 AI 回复 {reply_count} 条是否合理。
只回答"合理"或"不合理"。"""

        elif issue_type == "length":
            prompt = f"""你是一个审查员，需要判断 AI 的回复是否合理。

=== 聊天记录 ===
{context}

=== 当前情况 ===
用户发了 {msg_count} 条消息：
{user_message}

AI 有一条回复长度为 {details} 字（超过了 50 字限制）。

=== 判断标准 ===
"不合理"的定义：单条消息超过 50 字，且用户的问题并不需要长篇回复。

请根据聊天记录和当前情况，判断这条长消息是否合理。
只回答"合理"或"不合理"。"""
        else:
            return True
        
        resp = requests.post(
            api["url"],
            headers={"Authorization": f"Bearer {api['key']}", "Content-Type": "application/json"},
            json={"model": api["model"], "messages": [{"role": "user", "content": prompt}]},
            timeout=60
        )
        
        result = resp.json()
        if "choices" in result:
            answer = result["choices"][0]["message"]["content"].strip()
            print(f"[AI审查] 类型: {issue_type}, 用户消息数: {msg_count}, 结果: {answer}")
            return "合理" in answer
        else:
            print(f"[AI审查] API 返回无 choices: {result.get('error')}")
    except Exception as e:
        print(f"[AI审查] 出错: {e}")
    
    return True  # 出错默认合理

def evaluate_ai_response(user_id, user, current_channel, user_message, reply, msg_count):
    """
    评估 AI 回复，返回 (violations, need_rework)
    violations: 违规列表
    need_rework: 是否需要返工
    """
    print(f"[Debug] evaluate_ai_response 被调用: user_id={user_id}, msg_count={msg_count}")
    messages = [m.strip() for m in reply.split("|||") if m.strip()] if "|||" in reply else [reply.strip()]
    reply_count = len(messages)
    print(f"[Debug] 回复条数: {reply_count}, 限制: {msg_count * 3}")
    violations = []
    need_rework = False
    
    # 1. 检查分点列举（直接扣分）
    has_list, list_reason = check_reply_format_violation(reply)
    if has_list:
        old, new, rework = deduct_ai_points(user_id, list_reason)
        violations.append(f"分点列举: {list_reason}")
        if rework:
            need_rework = True
    
    # 2. 检查消息过长（审查）
    is_long, long_msg, length = check_messages_too_long(messages)
    if is_long:
        if not review_with_ai(user, current_channel, user_message, reply, msg_count, "length", str(length)):
            old, new, rework = deduct_ai_points(user_id, f"消息过长: {length}字")
            violations.append(f"消息过长: {length}字")
            if rework:
                need_rework = True
    
    # 3. 检查回复条数（审查）
    if reply_count > msg_count * 3:
        if not review_with_ai(user, current_channel, user_message, reply, msg_count, "count"):
            old, new, rework = deduct_ai_points(user_id, f"回复过多: {reply_count}条")
            violations.append(f"回复过多: {reply_count}条")
            if rework:
                need_rework = True
    
    # 4. 没有违规就加分
    if not violations:
        reward_ai_points(user_id)
    
    return violations, need_rework

# ========== 频道和记忆工具 ==========

def get_all_channels():
    try:
        resp = requests.get(
            "https://slack.com/api/conversations.list",
            headers={"Authorization": f"Bearer {SLACK_BOT_TOKEN}"},
            params={"types": "public_channel,private_channel", "limit": 200}
        )
        result = resp.json()
        if result.get("ok"):
            return [{"id": ch["id"], "name": ch["name"], "is_member": ch.get("is_member", False)} 
                    for ch in result.get("channels", [])]
    except:
        pass
    return []

def get_channel_list_for_ai():
    channels = get_all_channels()
    member_channels = [ch for ch in channels if ch.get("is_member")]
    return "、".join([f"#{ch['name']}" for ch in member_channels]) if member_channels else "（无）"

def load_all_memories():
    return jsonbin_load(JSONBIN_MEMORIES, {})

def save_all_memories(data):
    jsonbin_save(JSONBIN_MEMORIES, data)

def load_memories(user_id):
    return load_all_memories().get(user_id, [])

def save_memories(user_id, memories):
    all_mem = load_all_memories()
    all_mem[user_id] = memories
    save_all_memories(all_mem)

def add_memory(user_id, content):
    memories = load_memories(user_id)
    total = sum(len(m["content"]) for m in memories)
    while total + len(content) > MEMORY_LIMIT and memories:
        removed = memories.pop(0)
        total -= len(removed["content"])
    memories.append({"content": content, "time": get_time_str()})
    save_memories(user_id, memories)

def delete_memory(user_id, index):
    memories = load_memories(user_id)
    if 1 <= index <= len(memories):
        removed = memories.pop(index - 1)
        save_memories(user_id, memories)
        return removed["content"]
    return None

def clear_memories(user_id):
    save_memories(user_id, [])

def format_memories(user_id, show_numbers=True):
    memories = load_memories(user_id)
    if not memories:
        return ""
    return "\n".join([f"{i}. {m['content']}" if show_numbers else f"• {m['content']}" 
                      for i, m in enumerate(memories, 1)])

def get_channel_members(channel):
    try:
        resp = requests.get(
            "https://slack.com/api/conversations.members",
            headers={"Authorization": f"Bearer {SLACK_BOT_TOKEN}"},
            params={"channel": channel}
        )
        result = resp.json()
        return result.get("members", []) if result.get("ok") else []
    except:
        return []

def get_all_memories_for_channel(channel):
    members = get_channel_members(channel)
    parts = []
    for member_id in members:
        mem = format_memories(member_id, show_numbers=False)
        if mem:
            parts.append(f"【{get_display_name(member_id)}的记忆】\n{mem}")
    return "\n\n".join(parts)

def is_dm_channel(channel):
    return channel.startswith("D")

def get_channel_name(channel_id):
    if is_dm_channel(channel_id):
        return "私聊"
    try:
        resp = requests.get(
            "https://slack.com/api/conversations.info",
            headers={"Authorization": f"Bearer {SLACK_BOT_TOKEN}"},
            params={"channel": channel_id}
        )
        result = resp.json()
        if result.get("ok"):
            return "#" + result["channel"]["name"]
    except:
        pass
    return f"#{channel_id}"

def get_channel_name_only(channel_id):
    name = get_channel_name(channel_id)
    return name[1:] if name.startswith("#") else name

def get_channel_id_by_name(name):
    name = name.lower().strip().lstrip('#')
    channels = get_all_channels()
    for ch in channels:
        if ch["name"].lower() == name:
            return ch["id"]
    return None

def should_include_dm_history(user_id, channel):
    if is_dm_channel(channel):
        return True
    all_data = load_user_data()
    user = all_data.get(user_id, {})
    settings = user.get("channel_dm_settings", {})
    if channel in settings:
        return settings[channel]
    if get_channel_name_only(channel).lower() in [c.lower() for c in NO_DM_HISTORY_CHANNELS]:
        return False
    return True

def set_channel_dm_setting(user_id, channel, include_dm):
    all_data = load_user_data()
    if user_id not in all_data:
        all_data[user_id] = {}
    if "channel_dm_settings" not in all_data[user_id]:
        all_data[user_id]["channel_dm_settings"] = {}
    all_data[user_id]["channel_dm_settings"][channel] = include_dm
    save_user_data(all_data)

def get_username(user_id):
    try:
        resp = requests.get(
            "https://slack.com/api/users.info",
            headers={"Authorization": f"Bearer {SLACK_BOT_TOKEN}"},
            params={"user": user_id}
        )
        result = resp.json()
        if result.get("ok"):
            return result["user"]["name"]
    except:
        pass
    return user_id

def get_display_name(user_id):
    try:
        resp = requests.get(
            "https://slack.com/api/users.info",
            headers={"Authorization": f"Bearer {SLACK_BOT_TOKEN}"},
            params={"user": user_id}
        )
        result = resp.json()
        if result.get("ok"):
            return result["user"]["real_name"] or result["user"]["name"]
    except:
        pass
    return user_id

def get_user_dm_channel(user_id):
    try:
        resp = requests.post(
            "https://slack.com/api/conversations.open",
            headers={"Authorization": f"Bearer {SLACK_BOT_TOKEN}", "Content-Type": "application/json"},
            json={"users": user_id}
        )
        result = resp.json()
        if result.get("ok"):
            return result["channel"]["id"]
    except:
        pass
    return None

def is_unlimited_user(user_id):
    return get_username(user_id) in UNLIMITED_USERS

def check_and_use_points(user_id, api_name):
    if is_unlimited_user(user_id):
        return True, -1, None
    
    cost = APIS.get(api_name, {}).get("cost", 1)
    all_data = load_user_data()
    user = all_data.get(user_id, {})
    used = user.get("points_used", 0)
    remaining = POINTS_LIMIT - used
    
    if remaining < cost:
        return False, remaining, f"积分不足！剩余 {remaining}，需要 {cost}。"
    
    user["points_used"] = used + cost
    all_data[user_id] = user
    save_user_data(all_data)
    return True, POINTS_LIMIT - user["points_used"], None

def is_in_conversation(user_id, channel):
    all_data = load_user_data()
    user = all_data.get(user_id, {})
    last_active = user.get("channel_last_active", {}).get(channel, 0)
    return (get_cn_time().timestamp() - last_active) < CONVERSATION_TIMEOUT

def activate_channel_conversation(user_id, channel):
    """激活频道对话状态"""
    all_data = load_user_data()
    if user_id not in all_data:
        all_data[user_id] = {}
    if "channel_last_active" not in all_data[user_id]:
        all_data[user_id]["channel_last_active"] = {}
    all_data[user_id]["channel_last_active"][channel] = get_cn_time().timestamp()
    save_user_data(all_data)
    print(f"[Conversation] 激活用户 {user_id} 在频道 {channel} 的对话状态")

# ========== 历史记录构建 ==========

def build_history_messages(user, current_channel, api_name):
    max_tokens = API_TOKEN_LIMITS.get(api_name, 100000)
    available = int(max_tokens * 0.7)
    
    current_is_dm = is_dm_channel(current_channel)
    user_id = user.get("user_id", "")
    include_dm = should_include_dm_history(user_id, current_channel)
    
    all_msgs = []
    
    # 私聊历史
    if include_dm or current_is_dm:
        for m in user.get("dm_history", []):
            if m.get("content"):
                all_msgs.append({
                    "role": m["role"], "content": m["content"],
                    "timestamp": m.get("timestamp", 0),
                    "scene": "dm", "is_current": current_is_dm
                })
    
    # 频道历史
    if not current_is_dm:
        reset_time = user.get("channel_reset_times", {}).get(current_channel, 0)
        for m in get_channel_messages_since_reset(current_channel, reset_time):
            content = m.get("content", "")
            if not content:
                continue
            
            if m.get("is_bot"):
                role, formatted = "assistant", content
            elif m.get("user_id") == user_id:
                role, formatted = "user", content
            else:
                role, formatted = "user", f"[{m.get('username', '某人')}说] {content}"
            
            all_msgs.append({
                "role": role, "content": formatted,
                "timestamp": m.get("timestamp", 0),
                "scene": "channel", "is_current": True
            })
    
    all_msgs.sort(key=lambda x: x["timestamp"])
    
    total = sum(estimate_tokens(m["content"]) for m in all_msgs)
    while total > available and all_msgs:
        removed = all_msgs.pop(0)
        total -= estimate_tokens(removed["content"])
    
    result = []
    for m in all_msgs:
        content = m["content"]
        if not m["is_current"] and m["scene"] == "dm":
            content = f"[私聊] {content}"
        result.append({"role": m["role"], "content": content})
    
    return result

# ========== System Prompt ==========

def get_system_prompt(mode="long", user_id=None, channel=None, msg_count=1):
    memories_text = ""
    if channel:
        if is_dm_channel(channel) and user_id:
            mem = format_memories(user_id, show_numbers=False)
            if mem:
                memories_text = f"\n\n【{get_display_name(user_id)}的记忆】\n{mem}"
        else:
            mem = get_all_memories_for_channel(channel)
            if mem:
                memories_text = f"\n\n{mem}"

    current_scene = "私聊" if is_dm_channel(channel) else get_channel_name(channel)
    time_period, time_greeting = get_time_period()
    user_id_hint = f"\n当前用户 ID：{user_id}" if user_id else ""
    channel_list = get_channel_list_for_ai()
    
    base = f"""你是一个友好的AI助手。
当前时间: {get_time_str()}
现在是{time_period}（{time_greeting}）
当前场景：{current_scene}{user_id_hint}
可用频道：{channel_list}
{memories_text}

Slack 格式：*粗体* _斜体_ ~删除线~ `代码` ```代码块``` > 引用 <@用户ID>
禁止：# 标题、LaTeX、Markdown 表格

=== 场景意识 ===
- [私聊] 标签 = 私聊中的对话
- [某人说] 标签 = 频道里其他人说的话
- 私聊内容不要在频道里主动提起

*重要：你回复时绝对不要加这些标签！*
- 不要在回复开头加 [私聊]、[频道] 之类的标签
- 不要在回复开头加 # 号
- 这些标签是系统用来标记历史消息的，不是你该加的
- 你的回复直接写内容就好

=== 特殊能力 ===
[[定时|YYYY-MM-DD|HH:MM|内容]] - 定时消息
[[每日|HH:MM|主题]] - 每日消息
[[记忆|内容]] 或 [[记忆|用户ID|内容]] - 长期记忆
[[特殊日期|MM-DD|描述]] - 特殊日期
[[私聊|内容]] - 发私聊
[[发到频道|频道名|内容]] - 发到频道（频道名不要加#号，比如 [[发到频道|chat|你好]]）
[[反应|emoji]] - 表情反应

不需要回复时用：[不回]"""

    if mode == "short":
        points, status, points_prompt = get_ai_points_status(user_id) if user_id else (10, "ok", "")
        
        base += f"""

=== 短句模式 ===

像朋友发微信一样聊天。

*用户发了 {msg_count} 条消息*

{points_prompt}

*字数要求*：
- 最好 20 字以内
- 最多 50 字（超过要审查）
- 超过且不合理会扣分

*扣分规则*（仅短句模式生效）：
- 积分 > 0：每次扣 2 分
- 积分 ≤ 0：每次扣 5 分
- 积分到 -10：强制返工！

*会扣分的行为*：
1. 分点列举（1. 2. 3. 或 • 或 -）→ 直接扣分！
2. 单条 > 50 字 → 审查
3. 回复条数 > 用户消息数 × 3 → 审查

*格式*：
- 用 ||| 分隔多条消息
- 简短自然，像发微信

用户：在吗
你：在~

用户：今天好累
你：怎么啦|||工作太多了？"""

    return base

# ========== 解析隐藏命令 ==========

def parse_hidden_commands(reply, user_id, current_channel=None):
    schedules = load_schedules()
    if user_id not in schedules:
        schedules[user_id] = {"timed": [], "daily": [], "special_dates": {}}

    has_hidden = False
    original_reply = reply
    extra_actions = []

    # 定时消息
    for date_str, time_str, hint in re.findall(r'\[\[定时\|(\d{4}-\d{2}-\d{2})\|(\d{1,2}:\d{2})\|(.+?)\]\]', reply):
        h, m = time_str.split(":")
        schedules[user_id]["timed"].append({"date": date_str, "time": f"{int(h):02d}:{m}", "hint": hint})
        reply = reply.replace(f"[[定时|{date_str}|{time_str}|{hint}]]", "")
        has_hidden = True

    for time_str, hint in re.findall(r'\[\[定时\|(\d{1,2}:\d{2})\|([^\]]+?)\]\]', reply):
        h, m = time_str.split(":")
        schedules[user_id]["timed"].append({
            "date": get_cn_time().strftime("%Y-%m-%d"), "time": f"{int(h):02d}:{m}", "hint": hint
        })
        reply = reply.replace(f"[[定时|{time_str}|{hint}]]", "")
        has_hidden = True

    # 每日消息
    for time_str, topic in re.findall(r'\[\[每日\|(\d{1,2}:\d{2})\|(.+?)\]\]', reply):
        h, m = time_str.split(":")
        schedules[user_id]["daily"].append({"time": f"{int(h):02d}:{m}", "topic": topic})
        reply = reply.replace(f"[[每日|{time_str}|{topic}]]", "")
        has_hidden = True

    # 记忆
    for mem_uid, content in re.findall(r'\[\[记忆\|([A-Z0-9]+)\|(.+?)\]\]', reply):
        add_memory(mem_uid, content)
        reply = reply.replace(f"[[记忆|{mem_uid}|{content}]]", "")
        has_hidden = True

    for content in re.findall(r'\[\[记忆\|([^|]+?)\]\]', reply):
        if not re.match(r'^[A-Z0-9]+$', content):
            add_memory(user_id, content)
            reply = reply.replace(f"[[记忆|{content}]]", "")
            has_hidden = True

    # 特殊日期
    for date, desc in re.findall(r'\[\[特殊日期\|(\d{2}-\d{2})\|(.+?)\]\]', reply):
        schedules[user_id]["special_dates"][date] = desc
        reply = reply.replace(f"[[特殊日期|{date}|{desc}]]", "")
        has_hidden = True

    # 私聊
    for msg in re.findall(r'\[\[私聊\|(.+?)\]\]', reply):
        extra_actions.append({"type": "dm", "content": msg})
        reply = reply.replace(f"[[私聊|{msg}]]", "")
        has_hidden = True

    # 发到频道
    for ch, msg in re.findall(r'\[\[发到频道\|(\w+)\|(.+?)\]\]', reply):
        extra_actions.append({"type": "to_channel", "channel_name": ch, "content": msg})
        reply = reply.replace(f"[[发到频道|{ch}|{msg}]]", "")
        has_hidden = True

    # 反应
    for emoji in re.findall(r'\[\[反应\|(\w+)\]\]', reply):
        extra_actions.append({"type": "reaction", "emoji": emoji.lower()})
        reply = reply.replace(f"[[反应|{emoji}]]", "")
        has_hidden = True

    save_schedules(schedules)
    return re.sub(r'\n{3,}', '\n\n', reply).strip(), has_hidden, original_reply, extra_actions

# ========== API 调用 ==========

def call_ai(messages, api_name, has_image=False, max_retries=3):
    api = APIS.get(api_name, APIS[DEFAULT_API])
    
    if has_image and not api.get("vision"):
        return "当前模型不支持图片，请用 /model 切换。"

    for attempt in range(max_retries):
        try:
            resp = requests.post(
                api["url"],
                headers={"Authorization": f"Bearer {api['key']}", "Content-Type": "application/json"},
                json={"model": api["model"], "messages": messages},
                timeout=120
            )
            result = resp.json()
            
            if "choices" in result:
                return result["choices"][0]["message"]["content"]
            elif "error" in result:
                err = str(result.get("error", "")).lower()
                if any(x in err for x in ["upstream", "timeout", "do_request"]):
                    time.sleep(2 ** attempt)
                    continue
                return f"API 错误: {result.get('error')}"
        except requests.exceptions.Timeout:
            time.sleep(2 ** attempt)
        except Exception as e:
            if attempt == max_retries - 1:
                return f"出错了: {e}"
            time.sleep(2 ** attempt)
    
    return "API 请求失败 😢"

# ========== Slack 工具 ==========

def send_slack(channel, text):
    resp = requests.post(
        "https://slack.com/api/chat.postMessage",
        headers={"Authorization": f"Bearer {SLACK_BOT_TOKEN}"},
        json={"channel": channel, "text": text}
    )
    return resp.json().get("ts")

def update_slack(channel, ts, text):
    requests.post(
        "https://slack.com/api/chat.update",
        headers={"Authorization": f"Bearer {SLACK_BOT_TOKEN}"},
        json={"channel": channel, "ts": ts, "text": text}
    )

def delete_slack(channel, ts):
    requests.post(
        "https://slack.com/api/chat.delete",
        headers={"Authorization": f"Bearer {SLACK_BOT_TOKEN}"},
        json={"channel": channel, "ts": ts}
    )

def add_reaction(channel, ts, emoji):
    emoji = EMOJI_ALIASES.get(emoji.lower().replace(':', ''), emoji.lower())
    if emoji not in VALID_EMOJIS:
        return
    requests.post(
        "https://slack.com/api/reactions.add",
        headers={"Authorization": f"Bearer {SLACK_BOT_TOKEN}", "Content-Type": "application/json"},
        json={"channel": channel, "timestamp": ts, "name": emoji}
    )

def send_multiple_slack(channel, texts):
    for text in texts:
        if text.strip():
            send_slack(channel, text.strip())
            time.sleep(0.5)

def download_image(url):
    try:
        resp = requests.get(url, headers={"Authorization": f"Bearer {SLACK_BOT_TOKEN}"}, timeout=30)
        if resp.status_code == 200:
            return base64.b64encode(resp.content).decode('utf-8')
    except:
        pass
    return None

def execute_extra_actions(actions, user_id, channel, msg_ts=None, mode="long"):
    for action in actions:
        if action["type"] == "dm" and not is_dm_channel(channel):
            dm_ch = get_user_dm_channel(user_id)
            if dm_ch:
                content = action["content"]
                if mode == "short" and "|||" in content:
                    send_multiple_slack(dm_ch, content.split("|||"))
                else:
                    send_slack(dm_ch, content)
        
        elif action["type"] == "to_channel":
            target = get_channel_id_by_name(action["channel_name"])
            if target:
                content = action["content"]
                if mode == "short" and "|||" in content:
                    send_multiple_slack(target, content.split("|||"))
                else:
                    send_slack(target, content)
                
                # AI 发到频道后，激活该频道的对话状态
                activate_channel_conversation(user_id, target)
                # 记录到频道消息
                add_channel_message(target, "BOT", "AI", content, is_bot=True)
        
        elif action["type"] == "reaction" and msg_ts:
            add_reaction(channel, msg_ts, action["emoji"])

def check_pending_clear(user_id, channel):
    print(f"[PendingClear] 检查 {user_id}, logs={pending_clear_logs}")
    if user_id in pending_clear_logs:
        pending_clear_logs[user_id]["count"] -= 1
        remaining = pending_clear_logs[user_id]["count"]
        print(f"[PendingClear] 剩余 {remaining} 条")
        if remaining <= 0:
            print(f"[PendingClear] 开始清空")
            clear_user_chat_logs(user_id, pending_clear_logs[user_id].get("channel_only"))
            log_message(user_id, channel, None, None, is_reset=True)
            del pending_clear_logs[user_id]
            print(f"[PendingClear] 清空完成")

# ========== 频道观察 ==========

def should_trigger_observation(channel_id):
    global channel_message_counts
    channel_message_counts[channel_id] = channel_message_counts.get(channel_id, 0) + 1
    if channel_message_counts[channel_id] >= 10:
        channel_message_counts[channel_id] = 0
        return True
    return False

def observe_channel(channel_id):
    print(f"[Observe] 开始观察频道: {channel_id}")
    try:
        msgs = get_recent_channel_messages(channel_id, 10)
        print(f"[Observe] 获取到 {len(msgs)} 条消息")
        if not msgs:
            print(f"[Observe] 没有消息，退出")
            return
        
        text = "\n".join([
            f"[{'你' if m.get('is_bot') else m.get('username', '某人')}说] {m.get('content', '')}"
            for m in msgs
        ])
        
        members = get_channel_members(channel_id)
        if not members:
            return
        
        user_data_template = load_user_data().get(members[0], {})
        api = user_data_template.get("api", DEFAULT_API)
        
        prompt = f"""你正在观察频道 {get_channel_name(channel_id)}。
时间：{get_time_str()}

最近对话：
{text}

你可以：回复（直接写）、私聊某人（[[私聊给|用户名|内容]]）、不参与（[不回]）
不要强行参与。"""

        reply = call_ai([{"role": "user", "content": prompt}], api)
        
        if "[不回]" in reply or not reply.strip():
            return
        
        for username, content in re.findall(r'\[\[私聊给\|(.+?)\|(.+?)\]\]', reply):
            for m in msgs:
                if m.get("username") == username:
                    dm = get_user_dm_channel(m.get("user_id"))
                    if dm:
                        send_slack(dm, content)
                    break
            reply = reply.replace(f"[[私聊给|{username}|{content}]]", "")
        
        reply = reply.strip()
        if reply and "[不回]" not in reply:
            if "|||" in reply:
                send_multiple_slack(channel_id, reply.split("|||"))
            else:
                send_slack(channel_id, reply)
            add_channel_message(channel_id, "BOT", "AI", reply, is_bot=True)
            
            # AI 主动发言后，激活频道里所有成员的对话状态
            for member_id in members:
                activate_channel_conversation(member_id, channel_id)

    except Exception as e:
        print(f"[Observe] 出错: {e}")

# ========== 核心处理（带返工机制）==========

def process_message_with_rework(user_id, user, channel, text, api_name, mode, msg_count, typing_ts):
    """处理消息，如果积分到 -10 且违规则返工"""
    print(f"[Debug] process_message_with_rework: mode={mode}, msg_count={msg_count}")
    
    system = get_system_prompt(mode, user_id, channel, msg_count)
    messages = [{"role": "system", "content": system}]
    messages.extend(build_history_messages(user, channel, api_name))
    messages.append({"role": "user", "content": text})
    
    for attempt in range(MAX_REWORK_ATTEMPTS + 1):
        reply = call_ai(messages, api_name)
        visible, has_hidden, original, extra_actions = parse_hidden_commands(reply, user_id, channel)
        
        # 只在短句模式下评估
        if mode != "short":
            return visible, has_hidden, original, extra_actions, []
        
        violations, need_rework = evaluate_ai_response(user_id, user, channel, text, visible, msg_count)
        
        if not need_rework:
            return visible, has_hidden, original, extra_actions, violations
        
        if attempt >= MAX_REWORK_ATTEMPTS:
            print(f"[返工] 已达最大次数 {MAX_REWORK_ATTEMPTS}，使用默认回复")
            return "好的~", False, "好的~", [], violations
        
        print(f"[返工] 第 {attempt + 1} 次，违规: {violations}")
        
        # 添加返工提示
        rework_prompt = f"""
🚨 你的回复被拒绝了！违规内容：{', '.join(violations)}

你必须重新生成回复！要求：
- 每条消息 20 字以内（最多 50 字）
- 不要分点列举！
- 回复条数不要超过 {msg_count * 3} 条！
- 用户发了 {msg_count} 条消息

重新生成："""
        
        messages.append({"role": "assistant", "content": reply})
        messages.append({"role": "user", "content": rework_prompt})
    
    return "好的~", False, "好的~", [], violations

def process_message(user_id, channel, text, files=None, message_ts=None, msg_count=1):
    all_data = load_user_data()
    user = all_data.get(user_id, {
        "dm_history": [], "api": DEFAULT_API, "mode": "long",
        "points_used": 0, "user_id": user_id, "ai_points": AI_POINTS_DEFAULT
    })
    user["user_id"] = user_id

    api = user.get("api", DEFAULT_API)
    is_dm = is_dm_channel(channel)
    mode = user.get("mode", "long")

    can_use, remaining, msg = check_and_use_points(user_id, api)
    if not can_use:
        send_slack(channel, msg)
        return

    display_name = get_display_name(user_id)
    now = get_cn_time().timestamp()
    user["last_active"] = now
    
    if is_dm:
        user["dm_channel"] = channel
    else:
        user["last_channel"] = channel
        user.setdefault("channel_last_active", {})[channel] = now
        add_channel_message(channel, user_id, display_name, text)

    # 处理文件
    images, file_texts = [], []
    for f in (files or []):
        ftype, content = process_file(f)
        if ftype == "image":
            images.append(content)
        elif content:
            file_texts.append(content)
    
    full_text = (text + "\n\n" + "\n\n".join(file_texts)).strip() if file_texts else text

    log_message(user_id, channel, "user", full_text, username=display_name)
    
    # 保存用户数据
    all_data[user_id] = user
    save_user_data(all_data)

    typing_ts = send_slack(channel, "_Typing..._")
    
    # 如果有图片，构建特殊消息格式
    if images:
        system = get_system_prompt(mode, user_id, channel, msg_count)
        messages = [{"role": "system", "content": system}]
        messages.extend(build_history_messages(user, channel, api))
        
        content = []
        if full_text:
            content.append({"type": "text", "text": full_text})
        for url in images:
            img_data = download_image(url)
            if img_data:
                content.append({"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{img_data}"}})
        
        messages.append({"role": "user", "content": content})
        reply = call_ai(messages, api, has_image=True)
        visible, has_hidden, original, extra_actions = parse_hidden_commands(reply, user_id, channel)
        violations = []
    else:
        print(f"[Debug] 准备调用 process_message_with_rework, mode={mode}")
        visible, has_hidden, original, extra_actions, violations = process_message_with_rework(
            user_id, user, channel, full_text, api, mode, msg_count, typing_ts
        )

    model_name = APIS.get(api, {}).get("model", api)
    log_message(user_id, channel, "assistant", original, model=model_name, hidden=has_hidden)

    # 更新历史
    all_data = load_user_data()
    user = all_data.get(user_id, {})
    
    if is_dm:
        user.setdefault("dm_history", [])
        if full_text:
            user["dm_history"].append({"role": "user", "content": full_text, "timestamp": now})
        if original:
            user["dm_history"].append({"role": "assistant", "content": original, "timestamp": now + 0.001})
    else:
        if original:
            add_channel_message(channel, "BOT", "AI", original, is_bot=True)

    all_data[user_id] = user
    save_user_data(all_data)

    check_pending_clear(user_id, channel)
    execute_extra_actions(extra_actions, user_id, channel, message_ts, mode)

    if "[不回]" in visible or not visible.strip():
        delete_slack(channel, typing_ts)
    elif mode == "short" and "|||" in visible:
        parts = visible.split("|||")
        update_slack(channel, typing_ts, parts[0].strip())
        send_multiple_slack(channel, parts[1:])
    else:
        if remaining >= 0:
            visible += f"\n\n_剩余积分: {remaining}_"
        update_slack(channel, typing_ts, visible)

def delayed_process(user_id, channel, message_ts=None):
    time.sleep(5)
    
    if user_id not in pending_messages or not pending_messages[user_id]:
        return
    
    msgs = pending_messages[user_id]
    msg_count = len(msgs)
    combined = "\n".join(msgs)
    pending_messages[user_id] = []

    all_data = load_user_data()
    user = all_data.get(user_id, {
        "dm_history": [], "api": DEFAULT_API, "mode": "short",
        "points_used": 0, "user_id": user_id, "ai_points": AI_POINTS_DEFAULT
    })
    user["user_id"] = user_id

    api = user.get("api", DEFAULT_API)
    is_dm = is_dm_channel(channel)

    can_use, remaining, msg = check_and_use_points(user_id, api)
    if not can_use:
        send_slack(channel, msg)
        return

    typing_ts = send_slack(channel, "_Typing..._")
    display_name = get_display_name(user_id)
    log_message(user_id, channel, "user", combined, username=display_name)

    now = get_cn_time().timestamp()
    if is_dm:
        user["dm_channel"] = channel
    else:
        user["last_channel"] = channel
        user.setdefault("channel_last_active", {})[channel] = now
        add_channel_message(channel, user_id, display_name, combined)
    
    all_data[user_id] = user
    save_user_data(all_data)

    visible, has_hidden, original, extra_actions, violations = process_message_with_rework(
        user_id, user, channel, combined, api, "short", msg_count, typing_ts
    )

    model_name = APIS.get(api, {}).get("model", "未知")
    log_message(user_id, channel, "assistant", original, model=model_name, hidden=has_hidden)

    # 更新历史
    all_data = load_user_data()
    user = all_data.get(user_id, {})
    
    if is_dm:
        user.setdefault("dm_history", [])
        if combined:
            user["dm_history"].append({"role": "user", "content": combined, "timestamp": now})
        if original:
            user["dm_history"].append({"role": "assistant", "content": original, "timestamp": now + 0.001})
    else:
        if original:
            add_channel_message(channel, "BOT", "AI", original, is_bot=True)
    
    user["last_active"] = now
    all_data[user_id] = user
    save_user_data(all_data)

    check_pending_clear(user_id, channel)
    execute_extra_actions(extra_actions, user_id, channel, message_ts, "short")

    if "[不回]" in visible or not visible.strip():
        delete_slack(channel, typing_ts)
    elif "|||" in visible:
        parts = visible.split("|||")
        update_slack(channel, typing_ts, parts[0].strip())
        send_multiple_slack(channel, parts[1:])
    else:
        if remaining >= 0:
            visible += f"\n\n_剩余积分: {remaining}_"
        update_slack(channel, typing_ts, visible)

# ========== Slack 事件 ==========

@app.route("/slack/events", methods=["POST"])
def events():
    data = request.json

    if data.get("type") == "url_verification":
        return jsonify({"challenge": data["challenge"]})

    event_id = data.get("event_id")
    if event_id in processed_events:
        return jsonify({"ok": True})
    processed_events.add(event_id)
    if len(processed_events) > 1000:
        processed_events.clear()

    event = data.get("event", {})
    
    if event.get("type") not in ["app_mention", "message"]:
        return jsonify({"ok": True})
    # 过滤自己的消息（防止死循环），但允许其他bot（如清言）
    CLAUDE_BOT_IDS = ["B0AAJED3TEX", "B0AB9TGJBB3", "B056JA17RNU"]
    if event.get("bot_id") and event.get("bot_id") in CLAUDE_BOT_IDS:
        return jsonify({"ok": True})
    if event.get("subtype") and event.get("subtype") != "file_share":
        return jsonify({"ok": True})

    user_id = event.get("user")
    channel = event.get("channel")
    raw_text = event.get("text", "")
    text = re.sub(r'<@\w+>', '', raw_text).strip()
    message_ts = event.get("ts")

    # message_ts去重：防止app_mention和message.channels两个事件触发两次回复
    if message_ts in processed_events:
        return jsonify({"ok": True})
    processed_events.add(message_ts)
    if len(processed_events) > 1000:
        processed_events.clear()

    if message_ts in processed_file_events:
        return jsonify({"ok": True})
    
    files = event.get("files", [])
    if files:
        processed_file_events.add(message_ts)
        if len(processed_file_events) > 1000:
            processed_file_events.clear()

    if text.startswith("/"):
        return jsonify({"ok": True})

    is_dm = is_dm_channel(channel)
    is_mention = "<@" in raw_text
    in_conv = is_in_conversation(user_id, channel)
    
    # 非直接对话的频道消息
    if not is_dm and not is_mention and not in_conv:
        display_name = get_display_name(user_id)
        add_channel_message(channel, user_id, display_name, text)
        print(f"[Debug] 频道消息计数: {channel_message_counts.get(channel, 0)}")
        if should_trigger_observation(channel):
            print(f"[Debug] 触发频道观察: {channel}")
            threading.Thread(target=observe_channel, args=[channel]).start()
        return jsonify({"ok": True})
    
    if not text and not files:
        return jsonify({"ok": True})

    user_data = load_user_data().get(user_id, {})
    mode = user_data.get("mode", "long")

    if mode == "short" and not files:
        pending_messages.setdefault(user_id, []).append(text)
        if user_id in pending_timers:
            pending_timers[user_id].cancel()
        timer = threading.Timer(5.0, delayed_process, args=[user_id, channel, message_ts])
        timer.start()
        pending_timers[user_id] = timer
    else:
        threading.Thread(target=process_message, args=[user_id, channel, text, files, message_ts, 1]).start()

    return jsonify({"ok": True})

# ========== 斜杠命令 ==========

@app.route("/slack/commands", methods=["POST"])
def commands():
    cmd = request.form.get("command")
    user_id = request.form.get("user_id")
    channel = request.form.get("channel_id")
    text = request.form.get("text", "").strip()

    is_dm = is_dm_channel(channel)

    if cmd == "/reset":
        def do_reset():
            print(f"[Reset] 开始重置用户 {user_id}")
            data = load_user_data()
            if user_id in data:
                if is_dm:
                    data[user_id]["dm_history"] = []
                    data[user_id]["points_used"] = 0
                    scheds = load_schedules()
                    if user_id in scheds:
                        scheds[user_id] = {"timed": [], "daily": [], "special_dates": {}}
                        save_schedules(scheds)
                else:
                    data[user_id].setdefault("channel_reset_times", {})[channel] = get_cn_time().timestamp()
                save_user_data(data)
            print(f"[Reset] 重置完成")
        
        threading.Thread(target=do_reset).start()
        print(f"[Reset] 设置 pending_clear_logs: {user_id}")
        pending_clear_logs[user_id] = {"channel": channel, "count": 5, "channel_only": None if is_dm else channel}
        print(f"[Reset] pending_clear_logs 现在是: {pending_clear_logs}")
        
        msg = "✅ 已重置所有对话和定时任务！" if is_dm else f"✅ 已重置 {get_channel_name(channel)} 的对话！"
        return jsonify({"response_type": "in_channel", "text": f"{msg}\n📝 聊天记录将在 5 条消息后清空"})

    if cmd == "/memory":
        text_lower = text.lower()
        if not text_lower:
            mem = format_memories(user_id)
            total = sum(len(m["content"]) for m in load_memories(user_id))
            return jsonify({"response_type": "ephemeral", "text": f"📝 记忆（{total}/{MEMORY_LIMIT}字）：\n{mem}" if mem else "📝 暂无记忆"})
        
        if text_lower == "clear":
            threading.Thread(target=clear_memories, args=[user_id]).start()
            return jsonify({"response_type": "ephemeral", "text": "✅ 记忆已清空"})
        
        if text_lower.startswith("delete "):
            try:
                idx = int(text_lower[7:])
                removed = delete_memory(user_id, idx)
                return jsonify({"response_type": "ephemeral", "text": f"✅ 已删除: {removed}" if removed else "❌ 无效编号"})
            except:
                return jsonify({"response_type": "ephemeral", "text": "❌ 用法: /memory delete 编号"})
        
        return jsonify({"response_type": "ephemeral", "text": "用法:\n/memory - 查看\n/memory clear - 清空\n/memory delete 编号"})

    if cmd == "/model":
        all_data = load_user_data()
        
        if not text:
            info = "\n".join([f"{n} ({v['cost']}分) {'📷' if v.get('vision') else ''}" for n, v in APIS.items()])
            current = all_data.get(user_id, {}).get("api", DEFAULT_API)
            used = all_data.get(user_id, {}).get("points_used", 0)
            remaining = "∞" if is_unlimited_user(user_id) else f"{POINTS_LIMIT - used}/{POINTS_LIMIT}"
            return jsonify({"response_type": "ephemeral", "text": f"当前: {current}\n积分: {remaining}\n\n{info}"})
        
        if text in APIS:
            all_data.setdefault(user_id, {})["api"] = text
            save_user_data(all_data)
            v = APIS[text]
            return jsonify({"response_type": "ephemeral", "text": f"✅ {text} ({v['cost']}分，图片{'✅' if v.get('vision') else '❌'})"})
        
        return jsonify({"response_type": "ephemeral", "text": "❌ 无效模型"})

    if cmd == "/mode":
        all_data = load_user_data()
        text_lower = text.lower()
        
        if not text_lower:
            current = all_data.get(user_id, {}).get("mode", "long")
            return jsonify({"response_type": "ephemeral", "text": f"当前: {current}\n可用: long, short"})
        
        if text_lower in ["long", "short"]:
            all_data.setdefault(user_id, {})["mode"] = text_lower
            save_user_data(all_data)
            return jsonify({"response_type": "ephemeral", "text": f"✅ {text_lower}"})
        
        return jsonify({"response_type": "ephemeral", "text": "❌ 只能 long 或 short"})

    if cmd == "/dmhistory":
        if is_dm:
            return jsonify({"response_type": "ephemeral", "text": "❌ 只能在频道使用"})
        
        text_lower = text.lower()
        if not text_lower:
            include = should_include_dm_history(user_id, channel)
            return jsonify({"response_type": "ephemeral", "text": f"私聊记录: {'✅开启' if include else '❌关闭'}\n用法: /dm on|off"})
        
        if text_lower == "on":
            set_channel_dm_setting(user_id, channel, True)
            return jsonify({"response_type": "ephemeral", "text": "✅ 已开启私聊记录"})
        elif text_lower == "off":
            set_channel_dm_setting(user_id, channel, False)
            return jsonify({"response_type": "ephemeral", "text": "✅ 已关闭私聊记录"})
        
        return jsonify({"response_type": "ephemeral", "text": "❌ /dmhistory on 或 /dmhistory off"})

    if cmd == "/points":
        if is_unlimited_user(user_id):
            return jsonify({"response_type": "ephemeral", "text": "✨ 你是无限用户"})
        
        all_data = load_user_data()
        used = all_data.get(user_id, {}).get("points_used", 0)
        return jsonify({"response_type": "ephemeral", "text": f"剩余积分: {POINTS_LIMIT - used}/{POINTS_LIMIT}"})

    if cmd == "/aipoints":
        points = get_ai_points(user_id)
        status_msg = ""
        if points <= AI_POINTS_MIN:
            status_msg = "💀 最低分！再犯错要返工"
        elif points < 0:
            status_msg = "⚠️ 负分！每次扣5分"
        elif points <= 2:
            status_msg = "🚨 危险！"
        elif points <= 6:
            status_msg = "⚠️ 注意"
        else:
            status_msg = "👍 良好"
        
        return jsonify({"response_type": "ephemeral", "text": f"AI 积分: {points}/{AI_POINTS_MAX} {status_msg}"})

    return jsonify({"response_type": "ephemeral", "text": "未知命令"})

# ========== 定时任务 ==========

def run_scheduler():
    while True:
        try:
            now = get_cn_time()
            current_time = now.strftime("%H:%M")
            current_date = now.strftime("%Y-%m-%d")
            current_md = now.strftime("%m-%d")

            # 每日重置积分
            if current_time == "00:00":
                all_data = load_user_data()
                for uid in all_data:
                    all_data[uid]["points_used"] = 0
                save_user_data(all_data)
                print("[Scheduler] 用户积分已重置")

            all_data = load_user_data()
            schedules = load_schedules()
            print(f"[Scheduler] 加载 schedules: {len(schedules)} 个用户") # 新增调试打印

            # 如果 schedules 为空且不应该为空，跳过这次保存
            if not schedules: # 新增保护逻辑
                print(f"[Scheduler] schedules 为空，跳过本轮处理")
                time.sleep(60)
                continue

            for user_id, user in all_data.items():
                dm_channel = user.get("dm_channel")
                last_channel = user.get("last_channel")
                
                if not dm_channel and not last_channel:
                    dm_channel = get_user_dm_channel(user_id)
                    if dm_channel:
                        user["dm_channel"] = dm_channel
                        all_data[user_id] = user
                    else:
                        continue
                
                channel = dm_channel or last_channel
                if not channel:
                    continue

                user_scheds = schedules.get(user_id, {"timed": [], "daily": [], "special_dates": {}})
                api = user.get("api", DEFAULT_API)
                mode = user.get("mode", "long")

                # 定时消息
                new_timed = []
                for item in user_scheds.get("timed", []):
                    item_date = item.get("date", "")
                    item_time = item.get("time", "")
                    
                    if not item_date or not item_time:
                        continue
                    
                    # 标准化时间
                    if len(item_time.split(":")[0]) == 1:
                        item_time = "0" + item_time
                    
                    try:
                        target = datetime.strptime(f"{item_date} {item_time}", "%Y-%m-%d %H:%M")
                        target = target.replace(tzinfo=CN_TIMEZONE)
                    except:
                        new_timed.append(item)
                        continue
                    
                    if now >= target:
                        hint = item.get("hint", "")
                        print(f"[Scheduler] 触发定时: {hint[:30]}...")
                        
                        target_channel = dm_channel or channel
                        
                        system = get_system_prompt(mode, user_id, target_channel, 1)
                        system += f"\n\n=== 定时任务 ===\n你设定了：{hint}\n时间到了，发消息给用户。不想发就回复 [不发]"
                        
                        messages = [{"role": "system", "content": system}]
                        messages.extend(build_history_messages(user, target_channel, api))
                        messages.append({"role": "user", "content": "[定时任务触发]"}) # 新增占位符
                        
                        reply = call_ai(messages, api)
                        
                        if "[不发]" not in reply:
                            visible, has_hidden, original, extra = parse_hidden_commands(reply, user_id, target_channel)
                            
                            if visible.strip() and "[不回]" not in visible:
                                if mode == "short" and "|||" in visible:
                                    send_multiple_slack(target_channel, visible.split("|||"))
                                else:
                                    send_slack(target_channel, visible)
                                
                                log_message(user_id, target_channel, "assistant", original, 
                                           model=APIS.get(api, {}).get("model"), hidden=has_hidden)
                                
                                if is_dm_channel(target_channel):
                                    user.setdefault("dm_history", []).append({
                                        "role": "assistant", "content": original, 
                                        "timestamp": now.timestamp()
                                    })
                                else:
                                    add_channel_message(target_channel, "BOT", "AI", original, is_bot=True)
                                
                                execute_extra_actions(extra, user_id, target_channel, None, mode)
                    else:
                        new_timed.append(item)
                
                user_scheds["timed"] = new_timed

                # 每日消息
                for item in user_scheds.get("daily", []):
                    item_time = item.get("time", "")
                    if len(item_time.split(":")[0]) == 1:
                        item_time = "0" + item_time
                    
                    if item_time == current_time:
                        topic = item.get("topic", "")
                        print(f"[Scheduler] 触发每日: {topic[:30]}...")
                        
                        target_channel = dm_channel or channel
                        
                        system = get_system_prompt(mode, user_id, target_channel, 1)
                        system += f"\n\n=== 每日任务 ===\n主题：{topic}\n不想发就回复 [不发]"
                        
                        messages = [{"role": "system", "content": system}]
                        messages.extend(build_history_messages(user, target_channel, api))
                        messages.append({"role": "user", "content": "[定时任务触发]"}) # 新增占位符
                        
                        reply = call_ai(messages, api)
                        
                        if "[不发]" not in reply:
                            visible, has_hidden, original, extra = parse_hidden_commands(reply, user_id, target_channel)
                            
                            if visible.strip() and "[不回]" not in visible:
                                if mode == "short" and "|||" in visible:
                                    send_multiple_slack(target_channel, visible.split("|||"))
                                else:
                                    send_slack(target_channel, visible)
                                
                                log_message(user_id, target_channel, "assistant", original,
                                           model=APIS.get(api, {}).get("model"), hidden=has_hidden)
                                
                                if is_dm_channel(target_channel):
                                    user.setdefault("dm_history", []).append({
                                        "role": "assistant", "content": original,
                                        "timestamp": now.timestamp()
                                    })
                                else:
                                    add_channel_message(target_channel, "BOT", "AI", original, is_bot=True)
                                
                                execute_extra_actions(extra, user_id, target_channel, None, mode)

                # 特殊日期
                if current_time == "00:00":
                    special = user_scheds.get("special_dates", {}).get(current_md)
                    if special:
                        print(f"[Scheduler] 触发特殊日期: {special[:30]}...")
                        
                        target_channel = dm_channel or channel
                        
                        system = get_system_prompt(mode, user_id, target_channel, 1)
                        system += f"\n\n=== 特殊日期 ===\n今天是：{special}\n发一条祝福吧！不想发就回复 [不发]"
                        
                        messages = [{"role": "system", "content": system}]
                        messages.extend(build_history_messages(user, target_channel, api))
                        messages.append({"role": "user", "content": "[定时任务触发]"}) # 新增占位符
                        
                        reply = call_ai(messages, api)
                        
                        if "[不发]" not in reply:
                            visible, has_hidden, original, extra = parse_hidden_commands(reply, user_id, target_channel)
                            
                            if visible.strip() and "[不回]" not in visible:
                                if mode == "short" and "|||" in visible:
                                    send_multiple_slack(target_channel, visible.split("|||"))
                                else:
                                    send_slack(target_channel, visible)
                                
                                log_message(user_id, target_channel, "assistant", original,
                                           model=APIS.get(api, {}).get("model"), hidden=has_hidden)
                                
                                if is_dm_channel(target_channel):
                                    user.setdefault("dm_history", []).append({
                                        "role": "assistant", "content": original,
                                        "timestamp": now.timestamp()
                                    })
                                else:
                                    add_channel_message(target_channel, "BOT", "AI", original, is_bot=True)
                                
                                execute_extra_actions(extra, user_id, target_channel, None, mode)

                schedules[user_id] = user_scheds

            save_schedules(schedules)
            save_user_data(all_data)

        except Exception as e:
            print(f"[Scheduler] 出错: {e}")
            import traceback
            traceback.print_exc()

        time.sleep(60)

# ========== 启动 ==========

@app.route("/cron", methods=["GET", "POST"])
def cron_job():
    return jsonify({"ok": True, "message": "Using background scheduler"})

@app.route("/")
def home():
    return "Bot is running! 🤖"

scheduler_thread = threading.Thread(target=run_scheduler, daemon=True)
scheduler_thread.start()
print("[Startup] 定时任务线程已启动")

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 8080))
    app.run(host="0.0.0.0", port=port)
